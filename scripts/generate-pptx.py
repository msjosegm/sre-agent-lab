#!/usr/bin/env python3
"""Generate a clean, professional PowerPoint deck for the Azure SRE Agent MVP Lab Session."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BLUE = RGBColor(0x00, 0x33, 0x66)
AZURE_BLUE = RGBColor(0x00, 0x78, 0xD4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
ACCENT_GREEN = RGBColor(0x10, 0x7C, 0x10)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf

def add_para(tf, text, font_size=16, color=DARK_GRAY, bold=False, space_before=Pt(6)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_before = space_before
    return p

def add_shape_box(slide, left, top, width, height, fill_color, text="", font_size=14, font_color=WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(4)
    return shape

# ═══════════════════════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BLUE)
add_text_box(slide, 1, 1.5, 11, 1.5, "Azure SRE Agent", 54, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, 1, 3.2, 11, 1, "Your AI-Powered Operations Expert", 28, RGBColor(0xBB, 0xDD, 0xFF), False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 4.5, 11, 0.6, "Generally Available  —  March 10, 2026", 22, ACCENT_GREEN, True, PP_ALIGN.CENTER)
add_text_box(slide, 1, 5.8, 11, 0.6, "MVP Hands-On Lab Session", 20, RGBColor(0x99, 0xBB, 0xDD), False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 6.5, 11, 0.5, "sre.azure.com", 16, RGBColor(0x77, 0x99, 0xBB), False, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2: The Problem
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "Operations teams are drowning in toil", 36, DARK_BLUE, True)

# Pain points
pains = [
    ("2am alerts", "Manual runbook execution, error-prone under pressure"),
    ("5+ tool context-switching", "Logs → metrics → traces → code → wiki → Slack"),
    ("Knowledge silos", "Procedures locked in senior engineers' heads"),
    ("Triage backlogs", "Same classification work, every day, growing pile"),
]
y = 1.6
for title, desc in pains:
    add_shape_box(slide, 0.8, y, 0.4, 0.4, ACCENT_ORANGE)
    add_text_box(slide, 1.5, y, 4, 0.4, title, 20, DARK_GRAY, True)
    add_text_box(slide, 1.5, y + 0.4, 5, 0.4, desc, 14, RGBColor(0x66, 0x66, 0x66))
    y += 1.1

# Costs
add_text_box(slide, 7.5, 1.6, 5, 0.5, "The cost", 22, DARK_BLUE, True)
costs = [
    "MTTR measured in hours, not minutes",
    "Engineer burnout and attrition",
    "Customer satisfaction drops",
    "80% of response = gathering data",
]
tf = add_text_box(slide, 7.5, 2.3, 5, 3, "", 16, DARK_GRAY)
for c in costs:
    add_para(tf, f"→  {c}", 17, DARK_GRAY, space_before=Pt(12))

add_shape_box(slide, 2, 6.2, 9, 0.8, AZURE_BLUE, "What if an AI agent could handle the 80%?", 22, WHITE)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3: Introducing SRE Agent
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "Introducing Azure SRE Agent", 36, DARK_BLUE, True)
add_text_box(slide, 0.8, 1.2, 11, 0.6, "Your AI teammate that learns and grows with your team", 20, RGBColor(0x66, 0x66, 0x66))

# Three pillars
pillars = [
    ("Connect", "Azure resources\nIncident platforms\nAny API via MCP", AZURE_BLUE),
    ("Enhance", "Your runbooks\nYour architecture docs\nDomain subagents", ACCENT_GREEN),
    ("Achieve", "Faster investigations\nAutomated triage\nReliable remediation", ACCENT_ORANGE),
]
x = 1.2
for title, items, color in pillars:
    add_shape_box(slide, x, 2.2, 3.3, 0.7, color, title, 24, WHITE)
    add_text_box(slide, x, 3.1, 3.3, 2, items, 16, DARK_GRAY)
    x += 3.8

# Key differentiators
diffs = [
    "Not a chatbot — an autonomous agent that takes action",
    "Not generic AI — uses YOUR runbooks, YOUR data, YOUR procedures",
    "Not locked-in — connects to any tool via MCP",
]
y = 5.3
for d in diffs:
    add_text_box(slide, 1.2, y, 10, 0.4, f"✓  {d}", 17, DARK_GRAY, True)
    y += 0.45

# ═══════════════════════════════════════════════════════════════
# SLIDE 4: Four Superpowers
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "Four Superpowers", 36, DARK_BLUE, True)

powers = [
    ("Autonomous\nIncident Response", "Alert → investigate → remediate\nautonomously, 24/7", "Lower MTTR\nfrom hours to minutes", AZURE_BLUE),
    ("Lightning-Fast\nRoot Cause Analysis", "Correlates logs + metrics +\ntraces + code + deploys", "Actionable insights,\nnot data dumps", RGBColor(0x5B, 0x2C, 0x8B)),
    ("Extensible\nAutomation", "Connect any tool via MCP\nBuild self-healing workflows", "Eliminate toil,\nfree up engineers", ACCENT_GREEN),
    ("Conversational\nOperations", "Natural language + Python\nAsk questions, get evidence", "Anyone can\ninvestigate", ACCENT_ORANGE),
]

x = 0.5
for title, what, outcome, color in powers:
    add_shape_box(slide, x, 1.5, 3, 1.2, color, title, 16, WHITE)
    add_text_box(slide, x, 2.9, 3, 1.2, what, 13, DARK_GRAY)
    add_shape_box(slide, x, 4.3, 3, 0.9, LIGHT_GRAY, outcome, 13, DARK_GRAY)
    x += 3.2

add_text_box(slide, 0.8, 5.8, 11, 0.8, "Each superpower is demonstrated in today's lab", 18, AZURE_BLUE, True, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5: Three Personas
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "Three Personas — What You'll Try Today", 36, DARK_BLUE, True)

scenarios = [
    ("IT Operations / SRE",
     '"Stop running runbooks at 3am"',
     "Alert fires → agent auto-investigates\nusing your runbook → remediates\n→ creates GitHub issue with findings",
     "You wake up to results,\nnot raw alerts",
     AZURE_BLUE),
    ("Developer",
     '"From what happened to why + how to fix"',
     "Same investigation + source code\nsearch → finds exact file:line\n→ creates richer issue with fix",
     "Compare two GitHub issues\n— see the delta",
     RGBColor(0x5B, 0x2C, 0x8B)),
    ("Workflow Automation",
     '"Triage isn\'t the job, it\'s the tax"',
     "Agent triages GitHub issues:\nclassify → label → comment\nRuns on schedule, automatically",
     "Stop sorting tickets.\nStart shipping.",
     ACCENT_GREEN),
]

x = 0.5
for title, quote, desc, outcome, color in scenarios:
    add_shape_box(slide, x, 1.5, 3.8, 0.8, color, title, 18, WHITE)
    add_text_box(slide, x + 0.1, 2.5, 3.6, 0.5, quote, 13, RGBColor(0x66, 0x66, 0x66))
    add_text_box(slide, x + 0.1, 3.1, 3.6, 1.8, desc, 14, DARK_GRAY)
    add_shape_box(slide, x, 5.2, 3.8, 0.9, LIGHT_GRAY, outcome, 14, DARK_GRAY)
    x += 4.1

add_text_box(slide, 0.8, 6.4, 11, 0.5, "Core lab (IT Ops) works without GitHub  •  Developer + Workflow are bonus scenarios with GitHub", 14, RGBColor(0x66, 0x66, 0x66), False, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6: What azd up Deploys
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "One Command. Everything Configured.", 36, DARK_BLUE, True)

add_shape_box(slide, 3.5, 1.4, 6, 0.8, DARK_BLUE, "azd up", 32, WHITE)

components = [
    ("Infrastructure (Bicep)", ["SRE Agent", "Container Apps (API + Frontend)", "ACR, Log Analytics, App Insights", "Alert rules, Managed Identity + RBAC"]),
    ("Agent Config (REST APIs)", ["Knowledge base (2 runbooks)", "3 subagents (incident, code, triage)", "GitHub MCP connector", "Response plan, Scheduled task"]),
]

x = 1
for title, items in components:
    add_shape_box(slide, x, 2.6, 5.5, 0.6, AZURE_BLUE, title, 18, WHITE)
    y = 3.4
    for item in items:
        add_text_box(slide, x + 0.3, y, 5, 0.4, f"✓  {item}", 16, DARK_GRAY)
        y += 0.45
    x += 6

add_text_box(slide, 0.8, 5.8, 11, 0.6, "~8-12 minutes from clone to fully working agent", 20, ACCENT_GREEN, True, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7: GA Announcement
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_text_box(slide, 1, 0.5, 11, 1, "Generally Available", 44, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, 1, 1.5, 11, 0.6, "March 10, 2026", 28, ACCENT_GREEN, True, PP_ALIGN.CENTER)

features = [
    ("Enterprise SLA", "Production-grade support"),
    ("3 Regions", "East US 2, Sweden Central, Australia East"),
    ("Incident Platforms", "Azure Monitor, PagerDuty, ServiceNow"),
    ("MCP Extensibility", "Connect any API — GitHub, Datadog, Jira, yours"),
    ("Subagents", "Specialized agents for different domains"),
    ("Knowledge + Memory", "Your runbooks + learns from past incidents"),
]

y = 2.5
for title, desc in features:
    add_text_box(slide, 2, y, 4, 0.4, title, 20, WHITE, True)
    add_text_box(slide, 6, y, 6, 0.4, desc, 18, RGBColor(0xBB, 0xDD, 0xFF))
    y += 0.6

add_text_box(slide, 1, 6.3, 11, 0.5, "sre.azure.com/docs", 18, RGBColor(0x77, 0x99, 0xBB), False, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8: Before / After
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "The Transformation", 36, DARK_BLUE, True)

# Before
add_shape_box(slide, 0.8, 1.5, 5.5, 0.7, RGBColor(0xCC, 0x33, 0x33), "Before SRE Agent", 22, WHITE)
befores = [
    "2 hours every morning sorting tickets",
    "Context-switch between 5+ tools",
    "Knowledge locked in senior engineers",
    "Same diagnostic steps, error-prone at 3am",
]
y = 2.5
for b in befores:
    add_text_box(slide, 1.2, y, 5, 0.4, f"✗  {b}", 16, RGBColor(0xCC, 0x33, 0x33))
    y += 0.5

# After
add_shape_box(slide, 7, 1.5, 5.5, 0.7, ACCENT_GREEN, "After SRE Agent", 22, WHITE)
afters = [
    "Issues triaged before anyone logs in",
    "Investigations complete with evidence",
    "Runbooks executed consistently, 24/7",
    "Engineers solve problems, not gather data",
]
y = 2.5
for a in afters:
    add_text_box(slide, 7.4, y, 5, 0.4, f"✓  {a}", 16, ACCENT_GREEN)
    y += 0.5

add_shape_box(slide, 2, 5.2, 9, 1, DARK_BLUE, '"The best incident response is the one you sleep through."', 22, WHITE)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9: Let's Try It
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, AZURE_BLUE)
add_text_box(slide, 1, 1, 11, 1.5, "Let's Try It!", 54, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, 1, 3, 11, 1, "Open your lab instructions and follow along", 24, RGBColor(0xDD, 0xEE, 0xFF), False, PP_ALIGN.CENTER)

links = [
    ("Portal", "sre.azure.com"),
    ("Documentation", "sre.azure.com/docs"),
    ("Lab Repo", "github.com/msjosegm/sre-agent-lab"),
]
y = 4.5
for label, url in links:
    add_text_box(slide, 3, y, 3, 0.4, label, 20, WHITE, True, PP_ALIGN.RIGHT)
    add_text_box(slide, 6.5, y, 5, 0.4, url, 20, RGBColor(0xDD, 0xEE, 0xFF), False, PP_ALIGN.LEFT)
    y += 0.5

# ═══════════════════════════════════════════════════════════════
# SLIDE 10: Resources
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "Resources", 36, DARK_BLUE, True)

resources = [
    ("SRE Agent Portal", "sre.azure.com"),
    ("Documentation", "sre.azure.com/docs"),
    ("Lab Source Code", "github.com/msjosegm/sre-agent-lab"),
    ("Blog", "aka.ms/sreagent/blog"),
    ("Samples", "github.com/microsoft/sre-agent/tree/main/samples"),
]
y = 1.8
for label, url in resources:
    add_text_box(slide, 1.5, y, 4, 0.4, label, 20, DARK_GRAY, True)
    add_text_box(slide, 5.5, y, 7, 0.4, url, 18, AZURE_BLUE)
    y += 0.55

add_text_box(slide, 0.8, 5, 11, 0.6, "Blog posts to read after the lab:", 18, DARK_GRAY, True)
blogs = [
    "Stop Running Runbooks at 3am — Let Azure SRE Agent Do Your On-Call Grunt Work",
    "Extend SRE Agent with MCP — Build an Agentic Workflow to Triage Customer Issues",
]
y = 5.6
for b in blogs:
    add_text_box(slide, 1.2, y, 11, 0.4, f"→  {b}", 15, RGBColor(0x66, 0x66, 0x66))
    y += 0.45

# Save
output_path = "/tmp/sre-agent-lab/docs/Azure-SRE-Agent-MVP-Lab.pptx"
prs.save(output_path)
print(f"✅ PowerPoint saved to: {output_path}")
